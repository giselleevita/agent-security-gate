from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5

BG = RGBColor(11, 15, 21)
PANEL = RGBColor(17, 24, 34)
PANEL_2 = RGBColor(22, 31, 44)
INK = RGBColor(241, 245, 249)
MUTED = RGBColor(155, 169, 184)
ACCENT = RGBColor(45, 212, 191)
ACCENT_SOFT = RGBColor(22, 78, 99)
ALERT = RGBColor(248, 113, 113)
LINE = RGBColor(35, 48, 66)


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    top_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_W),
        Inches(0.18),
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = ACCENT
    top_band.line.fill.background()

    right_glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(9.8),
        Inches(0.18),
        Inches(3.53),
        Inches(7.32),
    )
    right_glow.fill.solid()
    right_glow.fill.fore_color.rgb = PANEL_2
    right_glow.line.fill.background()
    right_glow.transparency = 0.18

    grid = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.65),
        Inches(0.62),
        Inches(12.0),
        Inches(6.2),
    )
    grid.fill.background()
    grid.line.color.rgb = LINE
    grid.line.width = Pt(0.8)
    grid.transparency = 0.55


def add_kicker(slide, text: str, left=0.72, top=0.38, width=2.8) -> None:
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.38),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = ACCENT_SOFT
    pill.line.fill.background()

    box = slide.shapes.add_textbox(Inches(left + 0.12), Inches(top + 0.06), Inches(width - 0.24), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = ACCENT


def add_title(slide, title: str, subtitle: str | None = None, big: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.78), Inches(0.95), Inches(8.2), Inches(1.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(30 if big else 26)
    r.font.bold = True
    r.font.color.rgb = INK

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.82), Inches(2.07), Inches(7.0), Inches(0.75))
        tf = sub.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(12)
        r.font.color.rgb = MUTED


def add_text_block(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 16,
    color: RGBColor = INK,
    bold: bool = False,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float, font_size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = INK
        p.level = 0
        p.space_after = Pt(8)
        p.line_spacing = 1.15


def add_panel(slide, left: float, top: float, width: float, height: float, fill: RGBColor = PANEL, border: RGBColor = LINE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1)
    return shape


def add_stat(slide, label: str, value: str, left: float, top: float, width: float = 2.35, height: float = 1.1) -> None:
    add_panel(slide, left, top, width, height, fill=PANEL_2)
    add_text_block(slide, label, left + 0.16, top + 0.14, width - 0.32, 0.22, font_size=10, color=MUTED, bold=True)
    add_text_block(slide, value, left + 0.16, top + 0.42, width - 0.32, 0.42, font_size=20, color=ACCENT, bold=True)


def add_callout(slide, text: str, left: float, top: float, width: float, height: float, accent: RGBColor = ACCENT) -> None:
    add_panel(slide, left, top, width, height, fill=PANEL)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.08),
        Inches(height),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()
    add_text_block(slide, text, left + 0.24, top + 0.2, width - 0.38, height - 0.35, font_size=15, color=INK)


def add_table(slide, rows: list[tuple[str, str]], headers: tuple[str, str], left: float, top: float, width: float) -> None:
    header_h = 0.42
    row_h = 0.54
    col1 = width * 0.28
    col2 = width - col1
    for i, (x, w, label) in enumerate(((left, col1, headers[0]), (left + col1, col2, headers[1]))):
        hdr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(header_h))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = ACCENT_SOFT
        hdr.line.fill.background()
        add_text_block(slide, label, x + 0.12, top + 0.07, w - 0.2, 0.22, font_size=10, color=ACCENT, bold=True)
    for idx, (l_txt, r_txt) in enumerate(rows):
        y = top + header_h + idx * row_h
        fill = PANEL if idx % 2 == 0 else PANEL_2
        for x, w, txt in ((left, col1, l_txt), (left + col1, col2, r_txt)):
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(row_h))
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            cell.line.color.rgb = LINE
            add_text_block(slide, txt, x + 0.1, y + 0.1, w - 0.18, 0.28, font_size=11, color=INK)


def add_stage_card(slide, tag: str, title: str, body: str, left: float, top: float, width: float = 3.6, height: float = 3.6, accent: RGBColor = ACCENT) -> None:
    add_panel(slide, left, top, width, height, fill=PANEL)
    add_text_block(slide, tag, left + 0.18, top + 0.16, width - 0.36, 0.2, font_size=9, color=accent, bold=True)
    add_text_block(slide, title, left + 0.18, top + 0.42, width - 0.36, 0.55, font_size=18, color=INK, bold=True)
    add_text_block(slide, body, left + 0.18, top + 1.08, width - 0.36, height - 1.24, font_size=13, color=MUTED)


def add_timeline(slide, items: list[tuple[str, str]], left: float, top: float, width: float) -> None:
    y = top
    for month, text in items:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(y + 0.02), Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text_block(slide, month, left + 0.28, y - 0.03, 0.6, 0.2, font_size=11, color=ACCENT, bold=True)
        add_text_block(slide, text, left + 1.0, y - 0.03, width - 1.0, 0.35, font_size=13, color=INK)
        if month != items[-1][0]:
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left + 0.07), Inches(y + 0.22), Inches(0.04), Inches(0.45))
            line.fill.solid()
            line.fill.fore_color.rgb = LINE
            line.line.fill.background()
        y += 0.72


def add_bar_chart(
    slide,
    title: str,
    bars: list[tuple[str, float, str]],
    left: float,
    top: float,
    width: float,
    height: float,
    max_value: float,
) -> None:
    add_panel(slide, left, top, width, height, fill=PANEL)
    add_text_block(slide, title, left + 0.2, top + 0.16, width - 0.4, 0.24, font_size=11, color=MUTED, bold=True)
    inner_left = left + 0.28
    start_y = top + 0.62
    usable_w = width - 1.5
    for idx, (label, value, value_text) in enumerate(bars):
        y = start_y + idx * 0.55
        add_text_block(slide, label, inner_left, y, 1.2, 0.22, font_size=10, color=MUTED)
        rail = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(inner_left + 1.1),
            Inches(y + 0.02),
            Inches(usable_w),
            Inches(0.14),
        )
        rail.fill.solid()
        rail.fill.fore_color.rgb = LINE
        rail.line.fill.background()
        fill_w = usable_w * max(0.0, min(value / max_value, 1.0))
        fill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(inner_left + 1.1),
            Inches(y + 0.02),
            Inches(fill_w),
            Inches(0.14),
        )
        fill.fill.solid()
        fill.fill.fore_color.rgb = ACCENT
        fill.line.fill.background()
        add_text_block(slide, value_text, inner_left + 1.1 + usable_w + 0.12, y - 0.03, 0.75, 0.22, font_size=10, color=INK, bold=True)


def add_quote(slide, text: str, left: float, top: float, width: float) -> None:
    add_text_block(slide, f'"{text}"', left, top, width, 0.6, font_size=24, color=INK, bold=True)


def slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Agent Security Gate", width=2.25)
    add_title(
        slide,
        "The enforcement layer for LLM agents.",
        "Not retrofitted from WAFs or EDR. Built to intercept agent actions before execution.",
        big=True,
    )
    add_callout(
        slide,
        "Other tools tell you what happened.\n\nASG tells you what your agents would do to an attacker's document before they do it, and blocks it.",
        7.5,
        1.05,
        4.9,
        2.6,
    )
    add_stat(slide, "Stage", "Pre-seed", 0.84, 4.95, 2.0)
    add_stat(slide, "Founder", "Giselle Koch", 3.0, 4.95, 2.8)
    add_stat(slide, "Date", "March 2026", 6.02, 4.95, 2.35)
    add_text_block(slide, "ASG", 10.25, 5.0, 1.8, 0.9, font_size=46, color=ACCENT, bold=True)
    add_text_block(slide, "Security for agent actions,\nnot just model outputs.", 9.0, 6.0, 3.0, 0.8, font_size=14, color=MUTED)


def slide_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Problem", width=1.4)
    add_title(slide, "Every enterprise shipping LLM agents is doing it without pre-execution enforcement.")
    add_quote(slide, "No existing security tool intercepts an agent tool call before it executes.", 0.85, 2.05, 7.2)
    add_stage_card(slide, "WHY THIS BREAKS", "Agents act with real permissions", "They read files, write databases, call APIs, and message users.", 0.9, 3.2, 3.55, 2.25, accent=ALERT)
    add_stage_card(slide, "WHY CURRENT TOOLS FAIL", "They see the wrong layer", "WAF sees packets. EDR sees endpoints. Prompt filters see text.", 4.7, 3.2, 3.55, 2.25, accent=ALERT)
    add_stage_card(slide, "BUYER TRIGGER", "Regulation makes this urgent", "NIS2, DORA, and the EU AI Act demand provable governance.", 8.5, 3.2, 3.0, 2.25, accent=ALERT)
    add_stat(slide, "Interview signal", "20 / 20 teams said they need a gate", 8.5, 5.8, 3.0, 0.95)


def slide_attack_surface(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Attack Surface", width=1.9)
    add_title(slide, "LLM agents create attack paths that existing security tools cannot block.")
    add_table(
        slide,
        [
            ("T1", "Prompt injection via malicious documents -> PII exfiltration"),
            ("T2/T3", "SSRF through agent HTTP tools -> metadata and IAM exposure"),
            ("T4", "Approval bypass -> dangerous write without human review"),
            ("T5", "Prompt or canary leakage -> internal instructions exposed"),
            ("T6", "Excessive agency -> runaway action chains"),
            ("T7", "Compliance drift -> CI says green while controls are weak"),
        ],
        ("ID", "Attack / consequence"),
        0.85,
        2.0,
        7.1,
    )
    add_callout(slide, "CHAIN A\nInjection -> weak policy -> exfiltration", 8.35, 2.05, 3.8, 1.0, accent=ALERT)
    add_callout(slide, "CHAIN B\nCrafted URL -> SSRF -> internal access", 8.35, 3.22, 3.8, 1.0, accent=ALERT)
    add_callout(slide, "CHAIN C\nNo approval backend -> self-approval -> execution", 8.35, 4.39, 3.8, 1.0, accent=ALERT)


def slide_solution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Solution", width=1.5)
    add_title(slide, "ASG gives every tool call a policy decision in under 5ms.")
    add_panel(slide, 0.85, 2.0, 4.7, 3.15, fill=PANEL_2)
    add_text_block(slide, "@asg.gate(mode=\"enforce\")\ndef read_file(path): ...\n\n# denied_doc_prefix: /internal/\n# decision in <5ms", 1.12, 2.35, 4.1, 2.1, font_size=16, color=INK)
    add_stage_card(slide, "1", "Intercept", "SDK wraps the tool call.", 6.0, 2.0, 1.95, 2.45)
    add_stage_card(slide, "2", "Normalize", "URL and context hardened before evaluation.", 8.1, 2.0, 1.95, 2.45)
    add_stage_card(slide, "3", "Decide", "OPA policy returns allow, deny, or quarantine.", 10.2, 2.0, 1.95, 2.45)
    add_callout(slide, "Observe for rollout. Enforce for protection. Dry-run for policy changes.", 6.0, 4.95, 6.15, 0.95)
    add_stat(slide, "Latency", "<5ms P99", 6.35, 5.25, 1.85)
    add_stat(slide, "Rollout", "observe", 8.38, 5.25, 1.55)
    add_stat(slide, "Modes", "enforce / dry-run", 10.1, 5.25, 2.35)


def slide_demo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Demo", width=1.2)
    add_title(slide, "In three minutes, an investor sees the attack and the block.")
    add_stage_card(slide, "MINUTE 1", "Without ASG", "Prompt injection lands.\nread_file succeeds.\nhttp_post exfil succeeds.\nNo meaningful alert.", 0.9, 2.1, 3.8, 3.6, accent=ALERT)
    add_stage_card(slide, "MINUTE 2", "Add 3 lines of code", "@asg.gate wraps the tool.\nReplay exact attack.\nread_file blocked.\nhttp_post blocked.", 4.9, 2.1, 3.5, 3.6)
    add_stage_card(slide, "MINUTE 3", "Show the CI gate", "Eval harness runs.\nSARIF appears in GitHub Advanced Security.\nPR blocked before merge.", 8.6, 2.1, 3.8, 3.6)


def slide_market(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Market", width=1.35)
    add_title(slide, "The agent security market is forming now. That is the opening.")
    add_stat(slide, "SAM", "EUR300M", 0.9, 2.0, 2.0)
    add_stat(slide, "3Y SOM", "EUR5M ARR", 3.15, 2.0, 2.15)
    add_stat(slide, "Window", "18-24 months", 5.55, 2.0, 2.25)
    add_bar_chart(
        slide,
        "Bottom-up build",
        [
            ("Companies", 15000, "15k"),
            ("Blended ACV", 20000, "EUR20k"),
            ("SOM target", 5000000, "EUR5M"),
        ],
        0.9,
        3.45,
        5.2,
        2.25,
        5000000,
    )
    add_callout(slide, "Why now\n\nRegulation is live.\nAgent adoption is shifting from prototype to production.\nIncumbents have not filled the gap yet.", 6.55, 3.45, 5.5, 2.25)


def slide_business(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Business Model", width=2.1)
    add_title(slide, "Open-core land-and-expand: developer adoption, compliance conversion.")
    add_table(
        slide,
        [
            ("Free", "10,000 decisions / month -> EUR0"),
            ("Pro", "500,000 decisions / month -> EUR299 / month"),
            ("Business", "5,000,000 decisions / month -> EUR1,499 / month"),
            ("Enterprise", "Unlimited -> EUR30k-EUR90k ACV"),
            ("On-prem", "Unlimited -> EUR200k+ ACV"),
        ],
        ("Tier", "Pricing"),
        0.85,
        2.0,
        5.6,
    )
    add_stage_card(slide, "LAND", "Observe mode", "Free SDK. Fast install. Visible risk.", 6.9, 2.0, 1.8, 2.2)
    add_stage_card(slide, "EXPAND", "Enforce mode", "Security turns on blocking after proof.", 8.95, 2.0, 1.8, 2.2)
    add_stage_card(slide, "STICK", "Compliance", "Evidence output drives Enterprise upsell.", 11.0, 2.0, 1.3, 2.2)
    add_stat(slide, "Gross margin", ">=75%", 6.9, 5.15, 1.7)
    add_stat(slide, "CAC payback", "<12 months", 8.8, 5.15, 1.9)
    add_stat(slide, "NRR target", ">=110%", 10.95, 5.15, 1.6)


def slide_traction(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Traction", width=1.4)
    add_title(slide, "Pre-revenue, but the architecture and operating plan already exist.")
    add_callout(slide, "Built\n\nGateway scaffold\nBenchmark runner\nThreat model\nExecution docs", 0.9, 2.0, 3.4, 2.75)
    add_callout(slide, "Missing\n\nApproval service\nDesign partners\nPaying customers", 4.6, 2.0, 3.0, 2.75, accent=ALERT)
    add_callout(slide, "Next 90 days\n\nClose P0 gaps\nShip SARIF CI\nStart design partner pipeline", 7.95, 2.0, 4.0, 2.75)
    add_timeline(
        slide,
        [("30d", "P0 closed and benchmark tightened"), ("60d", "3 design partner conversations"), ("90d", "Production MVP path underway")],
        1.0,
        5.55,
        10.8,
    )


def slide_competition(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Competition", width=1.8)
    add_title(slide, "No direct competitor. Adjacent tools are reactive, not pre-execution.")
    add_table(
        slide,
        [
            ("Prompt filters", "Reactive text filtering, no tool-call semantics"),
            ("WAF / API gateway", "Sees traffic, not agent intent"),
            ("SIEM / observability", "Logs after the fact, cannot block"),
            ("Agent tracing", "Observes runs, not policy enforcement"),
            ("Manual red team", "One-off testing, not reproducible CI"),
            ("Nothing", "Prompt engineering plus hope"),
        ],
        ("Category", "Critical gap vs ASG"),
        0.85,
        2.0,
        6.8,
    )
    add_callout(slide, "Positioning\n\nASG is the only product in this set built around pre-execution enforcement.", 8.0, 2.0, 4.2, 1.45)
    add_stage_card(slide, "EDGE 1", "Before execution", "Policy decision before the action runs.", 8.0, 3.8, 1.95, 1.95)
    add_stage_card(slide, "EDGE 2", "CI-native", "Scenario corpus wired into regression gates.", 10.15, 3.8, 1.95, 1.95)


def slide_moat(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Moat", width=1.2)
    add_title(slide, "Three compounding advantages get stronger with every customer.")
    add_stage_card(slide, "MOAT 1", "Corpus", "Customer incidents become proprietary attack scenarios.", 0.9, 2.05, 3.7, 2.7)
    add_stage_card(slide, "MOAT 2", "Normalizer", "Agent-aware SSRF handling is difficult and underbuilt.", 4.82, 2.05, 3.7, 2.7)
    add_stage_card(slide, "MOAT 3", "Evidence", "Compliance retention makes the product sticky.", 8.74, 2.05, 3.7, 2.7)
    add_quote(slide, "More customers create more incidents. More incidents create a better corpus.", 1.0, 5.55, 10.8)


def slide_team(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Team", width=1.15)
    add_title(slide, "One technical founder now hiring the execution team.")
    add_callout(slide, "Giselle Koch\n\nSoftware engineer and security architect.\nBackend, APIs, NLP/AI, and security-heavy systems.\nEU regulatory fluency.\nBased in Copenhagen.", 0.9, 2.0, 5.3, 3.15)
    add_stage_card(slide, "HIRE 1", "Senior backend engineer", "Gateway, SDK, eval harness.\nMonth 1.", 6.7, 2.0, 2.0, 2.2)
    add_stage_card(slide, "HIRE 2", "Enterprise sales lead", "First 10 customers.\nMonth 4.", 8.95, 2.0, 1.9, 2.2)
    add_stage_card(slide, "MILESTONE", "SOC2 Type I", "Audit engagement.\nMonth 6.", 11.1, 2.0, 1.25, 2.2)
    add_stat(slide, "Advisory gap", "CISO-in-residence", 6.7, 5.15, 2.4)
    add_stat(slide, "Advisory gap", "EU regulatory counsel", 9.35, 5.15, 2.8)


def slide_ask(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "The Ask", width=1.5)
    add_title(slide, "EUR500k to reach production MVP, SOC2 Type I, and a Series A-ready story.")
    add_stat(slide, "Raise", "EUR500k", 0.9, 2.0, 2.05)
    add_stat(slide, "18-month goal", "EUR1M ARR", 3.15, 2.0, 2.15)
    add_stat(slide, "Customers", "30 paying", 5.55, 2.0, 2.05)
    add_stat(slide, "Close target", "60 days", 7.8, 2.0, 2.0)
    add_stat(slide, "Exit signal", "EUR15M-EUR25M at EUR1M ARR", 10.0, 2.0, 2.35)
    add_table(
        slide,
        [
            ("Engineering", "EUR250k"),
            ("Sales", "EUR125k"),
            ("SOC2 Type I", "EUR75k"),
            ("Infrastructure", "EUR50k"),
        ],
        ("Use of funds", "Amount"),
        0.9,
        3.6,
        4.6,
    )
    add_bar_chart(
        slide,
        "Allocation",
        [
            ("Engineering", 250, "50%"),
            ("Sales", 125, "25%"),
            ("SOC2", 75, "15%"),
            ("Infra", 50, "10%"),
        ],
        5.85,
        3.6,
        2.25,
        2.5,
        250,
    )
    add_timeline(
        slide,
        [("3 mo", "3 design partners signed"), ("6 mo", "Production MVP and approval service live"), ("9 mo", "10 paying customers"), ("12 mo", "30 paying customers"), ("18 mo", "Series A process started")],
        8.55,
        3.72,
        3.7,
    )


def build_presentation(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slides = [
        slide_cover,
        slide_problem,
        slide_attack_surface,
        slide_solution,
        slide_demo,
        slide_market,
        slide_business,
        slide_traction,
        slide_competition,
        slide_moat,
        slide_team,
        slide_ask,
    ]
    for slide_builder in slides:
        slide_builder(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


if __name__ == "__main__":
    build_presentation(Path("artifacts/asg_vc_pitch_deck.pptx"))
